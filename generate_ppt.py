"""Generate a professional PowerPoint presentation for the Movie Recommendation System project."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Constants & Color Palette
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Dark, premium color palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)        # Deep navy background
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)        # Card/section background
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)    # Bright sky blue accent
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)  # Soft purple accent
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)   # Emerald accent
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)   # Amber/yellow accent
ACCENT_ROSE = RGBColor(0xFB, 0x71, 0x85)    # Rose accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
MUTED_GRAY = RGBColor(0x94, 0xA3, 0xB8)
GRADIENT_START = RGBColor(0x1E, 0x3A, 0x5F)
GRADIENT_END = RGBColor(0x0F, 0x17, 0x2A)

OUTPUT_FILE = "Movie_Recommendation_System.pptx"

# ---------------------------------------------------------------------------
# Helper Functions
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color):
    """Set a solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None, corner_radius=None):
    """Add a rounded rectangle shape (card)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=15, color=LIGHT_GRAY,
                    bullet_color=ACCENT_BLUE, font_name="Calibri", spacing=Pt(8)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_icon_card(slide, left, top, w, h, icon_text, title, desc, accent_color):
    """Add a card with icon, title, and description."""
    card = add_shape(slide, left, top, w, h, BG_CARD, border_color=accent_color, border_width=1.5)
    # Icon circle
    circle_size = Inches(0.6)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.3), circle_size, circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    # Icon text
    add_text_box(slide, left + Inches(0.3), top + Inches(0.3), circle_size, circle_size,
                 icon_text, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, left + Inches(1.1), top + Inches(0.3), w - Inches(1.4), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
    # Description
    add_text_box(slide, left + Inches(1.1), top + Inches(0.7), w - Inches(1.4), h - Inches(1.0),
                 desc, font_size=12, color=MUTED_GRAY)


# ---------------------------------------------------------------------------
# Slide Builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    """Slide 1: Title Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, BG_DARK)

    # Decorative top accent bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_BLUE)

    # Title
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
                 "🎬 Movie Recommendation System", font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri")

    # Accent line under title
    add_accent_line(slide, Inches(4.5), Inches(2.9), Inches(4.3), ACCENT_BLUE, Pt(4))

    # Subtitle
    add_text_box(slide, Inches(2), Inches(3.2), Inches(9), Inches(0.7),
                 "A Hybrid Recommendation Engine with Interactive Visualizations",
                 font_size=22, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # Tech badges row
    badges = ["React + Vite", "FastAPI", "Scikit-learn", "D3.js", "Tailwind CSS"]
    badge_colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_PURPLE, ACCENT_AMBER, ACCENT_ROSE]
    badge_w = Inches(1.8)
    total_w = len(badges) * badge_w + (len(badges) - 1) * Inches(0.15)
    start_x = (SLIDE_WIDTH - total_w) / 2

    for i, (badge, bcolor) in enumerate(zip(badges, badge_colors)):
        x = start_x + i * (badge_w + Inches(0.15))
        card = add_shape(slide, x, Inches(4.3), badge_w, Inches(0.5), BG_CARD, border_color=bcolor, border_width=1)
        add_text_box(slide, x, Inches(4.33), badge_w, Inches(0.5),
                     badge, font_size=12, color=bcolor, bold=True, alignment=PP_ALIGN.CENTER)

    # Group project info
    add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5),
                 "Group Project  •  April 2026",
                 font_size=16, color=MUTED_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom decorative bar
    add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), ACCENT_PURPLE)


def slide_table_of_contents(prs):
    """Slide 2: Table of Contents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "📑 Table of Contents", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3), ACCENT_BLUE)

    sections = [
        ("01", "Project Overview", "Introduction and objectives", ACCENT_BLUE),
        ("02", "System Architecture", "Full-stack design and data flow", ACCENT_PURPLE),
        ("03", "Recommendation Algorithms", "Collaborative, SVD, and content-based methods", ACCENT_GREEN),
        ("04", "Tech Stack", "Frontend, backend, and ML libraries", ACCENT_AMBER),
        ("05", "Frontend Components", "React pages and interactive components", ACCENT_ROSE),
        ("06", "API Endpoints", "RESTful API design with FastAPI", ACCENT_BLUE),
        ("07", "Evaluation Metrics", "RMSE, MAE, Precision@K, Recall@K", ACCENT_PURPLE),
        ("08", "Results & Analysis", "Performance charts and findings", ACCENT_GREEN),
        ("09", "Key Features", "Search, graph, posters, and explanations", ACCENT_AMBER),
        ("10", "Conclusion & Future Work", "Summary and next steps", ACCENT_ROSE),
    ]

    col1_items = sections[:5]
    col2_items = sections[5:]

    for col_idx, items in enumerate([col1_items, col2_items]):
        x_base = Inches(0.8) + col_idx * Inches(6.3)
        for i, (num, title, desc, color) in enumerate(items):
            y = Inches(1.7) + i * Inches(1.05)
            # Number badge
            badge = add_shape(slide, x_base, y, Inches(0.55), Inches(0.55), color)
            add_text_box(slide, x_base, y + Inches(0.05), Inches(0.55), Inches(0.45),
                         num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
            # Title & desc
            add_text_box(slide, x_base + Inches(0.75), y, Inches(4.5), Inches(0.35),
                         title, font_size=16, color=WHITE, bold=True)
            add_text_box(slide, x_base + Inches(0.75), y + Inches(0.35), Inches(4.5), Inches(0.3),
                         desc, font_size=11, color=MUTED_GRAY)


def slide_project_overview(prs):
    """Slide 3: Project Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "01  Project Overview", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE)

    # Problem Statement Card
    add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(2.3), BG_CARD, border_color=ACCENT_BLUE, border_width=1)
    add_text_box(slide, Inches(1.1), Inches(1.85), Inches(5), Inches(0.4),
                 "🎯 Problem Statement", font_size=18, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.1), Inches(2.35), Inches(5.3), Inches(1.5),
                 "With thousands of movies available, users face choice overload. "
                 "Traditional browsing is inefficient. This project builds an intelligent "
                 "recommendation system that helps users discover movies they'll love, "
                 "powered by hybrid machine learning algorithms and interactive visualizations.",
                 font_size=14, color=LIGHT_GRAY)

    # Objectives Card
    add_shape(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(2.3), BG_CARD, border_color=ACCENT_GREEN, border_width=1)
    add_text_box(slide, Inches(7.3), Inches(1.85), Inches(5), Inches(0.4),
                 "🚀 Objectives", font_size=18, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, Inches(7.3), Inches(2.35), Inches(5), Inches(1.5), [
        "• Build a hybrid recommendation engine (Collaborative + SVD + Content)",
        "• Create an interactive React-based frontend with Tailwind CSS",
        "• Implement explainable AI – show WHY movies are recommended",
        "• Visualize movie similarity networks using D3.js force graphs",
    ], font_size=13, color=LIGHT_GRAY, spacing=Pt(6))

    # Key highlights row
    highlights = [
        ("🎬", "9,700+", "Movies in Dataset", ACCENT_BLUE),
        ("⭐", "100K+", "User Ratings", ACCENT_PURPLE),
        ("🤖", "3", "ML Algorithms", ACCENT_GREEN),
        ("📊", "4", "Evaluation Metrics", ACCENT_AMBER),
        ("🌐", "6", "API Endpoints", ACCENT_ROSE),
    ]
    card_w = Inches(2.2)
    gap = Inches(0.2)
    total = len(highlights) * card_w + (len(highlights) - 1) * gap
    sx = (SLIDE_WIDTH - total) / 2

    for i, (icon, value, label, color) in enumerate(highlights):
        x = sx + i * (card_w + gap)
        y = Inches(4.5)
        add_shape(slide, x, y, card_w, Inches(2.2), BG_CARD, border_color=color, border_width=1)
        add_text_box(slide, x, y + Inches(0.25), card_w, Inches(0.5),
                     icon, font_size=28, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(0.8), card_w, Inches(0.5),
                     value, font_size=30, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(1.4), card_w, Inches(0.5),
                     label, font_size=13, color=MUTED_GRAY, alignment=PP_ALIGN.CENTER)


def slide_architecture(prs):
    """Slide 4: System Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "02  System Architecture", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_PURPLE)

    # Architecture flow: User → Frontend → API → Backend → Data
    layers = [
        ("👤", "User", "Browser\nInteraction", ACCENT_BLUE, Inches(0.5)),
        ("⚛️", "React Frontend", "Vite + Tailwind\nD3 Graph Viz", ACCENT_PURPLE, Inches(2.8)),
        ("🔗", "REST API", "FastAPI\nEndpoints", ACCENT_GREEN, Inches(5.1)),
        ("🧠", "ML Engine", "Hybrid\nRecommender", ACCENT_AMBER, Inches(7.4)),
        ("📦", "Data Layer", "CSV Datasets\nPandas + Sklearn", ACCENT_ROSE, Inches(9.7)),
    ]

    y_center = Inches(2.8)
    box_w = Inches(2.0)
    box_h = Inches(2.5)

    for i, (icon, title, desc, color, x) in enumerate(layers):
        # Card
        add_shape(slide, x, y_center, box_w, box_h, BG_CARD, border_color=color, border_width=1.5)
        # Icon
        add_text_box(slide, x, y_center + Inches(0.2), box_w, Inches(0.5),
                     icon, font_size=30, alignment=PP_ALIGN.CENTER)
        # Title
        add_text_box(slide, x, y_center + Inches(0.8), box_w, Inches(0.4),
                     title, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        # Description
        add_text_box(slide, x + Inches(0.15), y_center + Inches(1.3), box_w - Inches(0.3), Inches(0.8),
                     desc, font_size=12, color=MUTED_GRAY, alignment=PP_ALIGN.CENTER)

        # Arrow between cards
        if i < len(layers) - 1:
            arrow_x = x + box_w + Inches(0.1)
            add_text_box(slide, arrow_x, y_center + Inches(0.9), Inches(0.5), Inches(0.5),
                         "→", font_size=28, color=MUTED_GRAY, alignment=PP_ALIGN.CENTER)

    # Data flow description
    add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), BG_CARD)
    add_text_box(slide, Inches(1.1), Inches(5.95), Inches(11), Inches(0.3),
                 "Data Flow", font_size=16, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.6),
                 "User searches a movie → React sends GET/POST request to FastAPI → "
                 "Backend loads pre-computed similarity matrices → Hybrid scoring combines "
                 "Collaborative (50%) + SVD (30%) + Content (20%) → Ranked results returned with explanations → "
                 "React renders cards, posters (via Wikipedia API), and D3 force graph.",
                 font_size=13, color=LIGHT_GRAY)


def slide_algorithms(prs):
    """Slide 5: Recommendation Algorithms."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "03  Recommendation Algorithms", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3), ACCENT_GREEN)

    algorithms = [
        {
            "title": "Collaborative Filtering",
            "weight": "50%",
            "color": ACCENT_BLUE,
            "icon": "👥",
            "desc": "Computes movie-movie cosine similarity from the user-movie rating matrix. "
                    "Movies liked by similar users are recommended.",
            "steps": [
                "1. Build user×movie rating matrix",
                "2. Transpose to get movie vectors over users",
                "3. Compute cosine similarity between all movies",
            ]
        },
        {
            "title": "SVD (Matrix Factorization)",
            "weight": "30%",
            "color": ACCENT_PURPLE,
            "icon": "🧮",
            "desc": "Uses TruncatedSVD to extract 50 latent features, then computes "
                    "cosine similarity on the compressed embeddings.",
            "steps": [
                "1. Transpose user-movie matrix (rows=movies)",
                "2. Apply TruncatedSVD (n_components=50)",
                "3. Cosine similarity on movie embeddings",
            ]
        },
        {
            "title": "Content-Based Filtering",
            "weight": "20%",
            "color": ACCENT_GREEN,
            "icon": "🎭",
            "desc": "Multi-label binarizes movie genres and computes cosine similarity "
                    "on genre vectors. Recommends movies with similar genres.",
            "steps": [
                "1. Parse genres into multi-label sets",
                "2. Binarize with MultiLabelBinarizer",
                "3. Cosine similarity on genre vectors",
            ]
        },
    ]

    card_w = Inches(3.8)
    gap = Inches(0.25)
    total = len(algorithms) * card_w + (len(algorithms) - 1) * gap
    sx = (SLIDE_WIDTH - total) / 2

    for i, algo in enumerate(algorithms):
        x = sx + i * (card_w + gap)
        y = Inches(1.7)
        h = Inches(5.3)
        add_shape(slide, x, y, card_w, h, BG_CARD, border_color=algo["color"], border_width=1.5)

        # Icon + Title
        add_text_box(slide, x, y + Inches(0.2), card_w, Inches(0.5),
                     algo["icon"], font_size=32, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(0.8), card_w, Inches(0.4),
                     algo["title"], font_size=17, color=algo["color"], bold=True, alignment=PP_ALIGN.CENTER)

        # Weight badge
        badge_w = Inches(1.0)
        add_shape(slide, x + (card_w - badge_w) / 2, y + Inches(1.3), badge_w, Inches(0.35),
                  algo["color"])
        add_text_box(slide, x + (card_w - badge_w) / 2, y + Inches(1.3), badge_w, Inches(0.35),
                     f"Weight: {algo['weight']}", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x + Inches(0.25), y + Inches(1.9), card_w - Inches(0.5), Inches(1.2),
                     algo["desc"], font_size=12, color=LIGHT_GRAY)

        # Steps
        add_text_box(slide, x + Inches(0.25), y + Inches(3.2), card_w - Inches(0.5), Inches(0.3),
                     "How it works:", font_size=13, color=algo["color"], bold=True)
        add_bullet_list(slide, x + Inches(0.25), y + Inches(3.6), card_w - Inches(0.5), Inches(1.5),
                        algo["steps"], font_size=11, color=MUTED_GRAY, spacing=Pt(5))

    # Hybrid formula at bottom
    add_shape(slide, Inches(2), Inches(7.15), Inches(9.3), Inches(0.25), BG_CARD)
    add_text_box(slide, Inches(2), Inches(7.15), Inches(9.3), Inches(0.25),
                 "Hybrid Score = 0.5 × Collaborative + 0.3 × SVD + 0.2 × Content",
                 font_size=13, color=ACCENT_AMBER, bold=True, alignment=PP_ALIGN.CENTER)


def slide_tech_stack(prs):
    """Slide 6: Tech Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "04  Tech Stack", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_AMBER)

    # Frontend card
    add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.3), BG_CARD, border_color=ACCENT_BLUE, border_width=1.5)
    add_text_box(slide, Inches(1.2), Inches(1.9), Inches(5), Inches(0.4),
                 "⚛️  Frontend", font_size=22, color=ACCENT_BLUE, bold=True)

    frontend_items = [
        ("React (Vite)", "Fast dev server, HMR, optimized builds"),
        ("Tailwind CSS", "Utility-first responsive styling"),
        ("react-force-graph-2d", "D3-powered interactive similarity graphs"),
        ("Axios", "HTTP client for API communication"),
        ("React Router", "Client-side routing for SPA navigation"),
    ]
    for i, (tech, desc) in enumerate(frontend_items):
        y = Inches(2.5) + i * Inches(0.8)
        add_text_box(slide, Inches(1.4), y, Inches(2.2), Inches(0.35),
                     f"▸ {tech}", font_size=14, color=WHITE, bold=True)
        add_text_box(slide, Inches(3.6), y, Inches(2.8), Inches(0.35),
                     desc, font_size=12, color=MUTED_GRAY)

    # Backend card
    add_shape(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(5.3), BG_CARD, border_color=ACCENT_GREEN, border_width=1.5)
    add_text_box(slide, Inches(7.4), Inches(1.9), Inches(5), Inches(0.4),
                 "🐍  Backend & ML", font_size=22, color=ACCENT_GREEN, bold=True)

    backend_items = [
        ("FastAPI", "High-performance async Python web framework"),
        ("Pandas", "Data manipulation and analysis"),
        ("Scikit-learn", "Cosine similarity, SVD, ML binarizers"),
        ("httpx", "Async HTTP client for Wikipedia poster API"),
        ("Pydantic", "Data validation and API schemas"),
    ]
    for i, (tech, desc) in enumerate(backend_items):
        y = Inches(2.5) + i * Inches(0.8)
        add_text_box(slide, Inches(7.6), y, Inches(2.2), Inches(0.35),
                     f"▸ {tech}", font_size=14, color=WHITE, bold=True)
        add_text_box(slide, Inches(9.8), y, Inches(2.5), Inches(0.35),
                     desc, font_size=12, color=MUTED_GRAY)


def slide_frontend_components(prs):
    """Slide 7: Frontend Components."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "05  Frontend Components", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3), ACCENT_ROSE)

    components = [
        {
            "name": "Home.jsx",
            "type": "Page",
            "color": ACCENT_BLUE,
            "icon": "🏠",
            "desc": "Landing page with search bar, featured movies, and navigation to all sections."
        },
        {
            "name": "Recommendations.jsx",
            "type": "Page",
            "color": ACCENT_PURPLE,
            "icon": "🎯",
            "desc": "Displays top-10 recommended movies for a selected title with similarity scores & explanations."
        },
        {
            "name": "PersonalizedPage.jsx",
            "type": "Page",
            "color": ACCENT_GREEN,
            "icon": "💡",
            "desc": "Multi-movie input for personalized recommendations. Users select multiple liked movies."
        },
        {
            "name": "Popular.jsx",
            "type": "Page",
            "color": ACCENT_AMBER,
            "icon": "🔥",
            "desc": "Top-20 highest-rated movies with minimum 50 ratings. Shows average scores and genres."
        },
        {
            "name": "MovieCard.jsx",
            "type": "Component",
            "color": ACCENT_ROSE,
            "icon": "🎬",
            "desc": "Reusable card with poster (Wikipedia API), title, genre tags, score badge & explanations."
        },
        {
            "name": "MovieGraph.jsx",
            "type": "Component",
            "color": ACCENT_BLUE,
            "icon": "🕸️",
            "desc": "Interactive D3 force-directed graph showing movie similarity network with weighted edges."
        },
        {
            "name": "SearchBar.jsx",
            "type": "Component",
            "color": ACCENT_PURPLE,
            "icon": "🔍",
            "desc": "Auto-completing search with debounced API calls. Filters movies as the user types."
        },
        {
            "name": "Navbar.jsx",
            "type": "Component",
            "color": ACCENT_GREEN,
            "icon": "📌",
            "desc": "Responsive navigation bar with links to Home, Popular, and Personalized pages."
        },
    ]

    cols = 4
    rows = 2
    card_w = Inches(2.8)
    card_h = Inches(2.5)
    gap_x = Inches(0.25)
    gap_y = Inches(0.25)
    total_w = cols * card_w + (cols - 1) * gap_x
    sx = (SLIDE_WIDTH - total_w) / 2

    for idx, comp in enumerate(components):
        col = idx % cols
        row = idx // cols
        x = sx + col * (card_w + gap_x)
        y = Inches(1.7) + row * (card_h + gap_y)

        add_shape(slide, x, y, card_w, card_h, BG_CARD, border_color=comp["color"], border_width=1)
        # Icon + name
        add_text_box(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), Inches(0.4),
                     f"{comp['icon']}  {comp['name']}", font_size=13, color=comp["color"], bold=True)
        # Type badge
        badge_w = Inches(0.9)
        add_shape(slide, x + Inches(0.15), y + Inches(0.6), badge_w, Inches(0.25), comp["color"])
        add_text_box(slide, x + Inches(0.15), y + Inches(0.6), badge_w, Inches(0.25),
                     comp["type"], font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Description
        add_text_box(slide, x + Inches(0.15), y + Inches(1.0), card_w - Inches(0.3), Inches(1.3),
                     comp["desc"], font_size=11, color=MUTED_GRAY)


def slide_api_endpoints(prs):
    """Slide 8: API Endpoints."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "06  API Endpoints", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE)

    endpoints = [
        ("GET", "/movies?q=...", "Search & list movies", "Returns up to 50 matching movies with title, genres, and movieId"),
        ("GET", "/recommend/{title}", "Get recommendations", "Top-10 similar movies with hybrid scores and explanations"),
        ("GET", "/popular", "Popular movies", "Top-20 highest-rated movies (min 50 ratings)"),
        ("POST", "/personalized-recommendations", "Personalized picks", "Recommendations based on multiple liked movies"),
        ("GET", "/graph/{title}", "Similarity graph", "Nodes & weighted links for D3 force-graph visualization"),
        ("GET", "/poster/{title}", "Movie poster", "Wikipedia thumbnail URL for movie card display"),
        ("GET", "/metrics", "Evaluation metrics", "RMSE, MAE, Precision@10, Recall@10 performance data"),
    ]

    # Table header
    header_y = Inches(1.6)
    add_shape(slide, Inches(0.8), header_y, Inches(11.7), Inches(0.5), ACCENT_BLUE)
    cols_x = [Inches(0.8), Inches(1.8), Inches(4.5), Inches(7.0)]
    cols_w = [Inches(1.0), Inches(2.7), Inches(2.5), Inches(5.5)]
    headers = ["Method", "Endpoint", "Purpose", "Description"]

    for cx, cw, h in zip(cols_x, cols_w, headers):
        add_text_box(slide, cx, header_y + Inches(0.05), cw, Inches(0.4),
                     h, font_size=13, color=WHITE, bold=True)

    # Table rows
    for i, (method, endpoint, purpose, desc) in enumerate(endpoints):
        y = Inches(2.2) + i * Inches(0.72)
        bg = BG_CARD if i % 2 == 0 else BG_DARK
        add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.65), bg)

        method_color = ACCENT_GREEN if method == "GET" else ACCENT_AMBER
        # Method badge
        add_shape(slide, Inches(0.9), y + Inches(0.12), Inches(0.7), Inches(0.35), method_color)
        add_text_box(slide, Inches(0.9), y + Inches(0.12), Inches(0.7), Inches(0.35),
                     method, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, Inches(1.8), y + Inches(0.1), Inches(2.7), Inches(0.45),
                     endpoint, font_size=11, color=ACCENT_BLUE, bold=True)
        add_text_box(slide, Inches(4.5), y + Inches(0.1), Inches(2.5), Inches(0.45),
                     purpose, font_size=12, color=WHITE)
        add_text_box(slide, Inches(7.0), y + Inches(0.1), Inches(5.3), Inches(0.45),
                     desc, font_size=11, color=MUTED_GRAY)


def slide_evaluation_metrics(prs):
    """Slide 9: Evaluation Metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "07  Evaluation Metrics", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3), ACCENT_PURPLE)

    # Methodology card
    add_shape(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.5), BG_CARD, border_color=ACCENT_PURPLE, border_width=1)
    add_text_box(slide, Inches(1.1), Inches(1.75), Inches(11), Inches(0.35),
                 "📐 Methodology", font_size=18, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(1.1), Inches(2.15), Inches(11), Inches(0.8),
                 "The system uses an 80/20 train-test split on user ratings. TruncatedSVD (50 components) "
                 "creates a low-rank approximation of the user-movie matrix as predicted ratings. "
                 "A \"relevant\" item is defined as a movie the user rated ≥ 4.0 in the test set.",
                 font_size=13, color=LIGHT_GRAY)

    # Metric cards
    metrics = [
        ("RMSE", "Root Mean Squared Error", "√(mean((pred - actual)²))", "Penalizes larger errors more heavily", ACCENT_BLUE),
        ("MAE", "Mean Absolute Error", "mean(|pred - actual|)", "Average magnitude of prediction errors", ACCENT_GREEN),
        ("P@10", "Precision at 10", "|relevant ∩ top-10| / 10", "Fraction of top-10 that are relevant", ACCENT_AMBER),
        ("R@10", "Recall at 10", "|relevant ∩ top-10| / |relevant|", "How many relevant items are captured", ACCENT_ROSE),
    ]

    card_w = Inches(2.7)
    gap = Inches(0.25)
    total = len(metrics) * card_w + (len(metrics) - 1) * gap
    sx = (SLIDE_WIDTH - total) / 2

    for i, (short, full, formula, desc, color) in enumerate(metrics):
        x = sx + i * (card_w + gap)
        y = Inches(3.5)
        h = Inches(3.5)
        add_shape(slide, x, y, card_w, h, BG_CARD, border_color=color, border_width=1.5)

        # Metric name
        add_text_box(slide, x, y + Inches(0.25), card_w, Inches(0.5),
                     short, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(0.8), card_w, Inches(0.35),
                     full, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Formula box
        add_shape(slide, x + Inches(0.2), y + Inches(1.3), card_w - Inches(0.4), Inches(0.5), BG_DARK)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.35), card_w - Inches(0.4), Inches(0.4),
                     formula, font_size=10, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x + Inches(0.2), y + Inches(2.1), card_w - Inches(0.4), Inches(1.0),
                     desc, font_size=12, color=MUTED_GRAY, alignment=PP_ALIGN.CENTER)


def slide_results(prs):
    """Slide 10: Results & Analysis – embed the chart images."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "08  Results & Analysis", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_GREEN)

    # Try to embed actual chart images
    chart_files = [
        ("backend/error_metrics.png", "Error Metrics (RMSE & MAE)"),
        ("backend/retrieval_metrics.png", "Retrieval Metrics (P@10 & R@10)"),
        ("backend/top_genres.png", "Top Genre Distribution"),
        ("backend/rating_distribution.png", "Rating Distribution"),
    ]

    positions = [
        (Inches(0.5), Inches(1.5), Inches(5.9), Inches(2.7)),
        (Inches(6.8), Inches(1.5), Inches(5.9), Inches(2.7)),
        (Inches(0.5), Inches(4.5), Inches(5.9), Inches(2.7)),
        (Inches(6.8), Inches(4.5), Inches(5.9), Inches(2.7)),
    ]

    for (img_path, label), (x, y, w, h) in zip(chart_files, positions):
        full_path = os.path.join("d:/movie-app", img_path)
        # Card background
        add_shape(slide, x, y, w, h, BG_CARD, border_color=MUTED_GRAY, border_width=0.5)
        if os.path.exists(full_path):
            # Leave room for label
            img_h = h - Inches(0.4)
            slide.shapes.add_picture(full_path, x + Inches(0.1), y + Inches(0.1),
                                     w - Inches(0.2), img_h - Inches(0.1))
            add_text_box(slide, x, y + img_h, w, Inches(0.35),
                         label, font_size=11, color=MUTED_GRAY, alignment=PP_ALIGN.CENTER)
        else:
            add_text_box(slide, x, y + Inches(1.0), w, Inches(0.5),
                         f"[{label}]", font_size=14, color=MUTED_GRAY, alignment=PP_ALIGN.CENTER)


def slide_key_features(prs):
    """Slide 11: Key Features."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "09  Key Features", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_AMBER)

    features = [
        ("🔍", "Smart Search", "Auto-completing search with debounced API calls. "
         "Filters movies as users type, returning up to 50 results.", ACCENT_BLUE),
        ("🤖", "Explainable AI", "Each recommendation includes human-readable explanations: "
         "shared genres, collaborative signals, and latent feature patterns.", ACCENT_PURPLE),
        ("🕸️", "Similarity Graph", "Interactive D3 force-directed network visualization. "
         "Nodes represent movies, edges show similarity weights.", ACCENT_GREEN),
        ("🖼️", "Auto Posters", "Movie posters are fetched automatically from Wikipedia's "
         "REST API. No API key needed. Results are LRU-cached (2000 entries).", ACCENT_AMBER),
        ("💡", "Personalized Picks", "Input multiple favorite movies to get a tailored "
         "recommendation list based on averaged hybrid similarity vectors.", ACCENT_ROSE),
        ("📊", "Metrics Dashboard", "View real-time evaluation metrics: RMSE, MAE, "
         "Precision@10, and Recall@10 from an 80/20 train-test split.", ACCENT_BLUE),
    ]

    cols = 3
    rows = 2
    card_w = Inches(3.7)
    card_h = Inches(2.3)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)
    total_w = cols * card_w + (cols - 1) * gap_x
    sx = (SLIDE_WIDTH - total_w) / 2

    for idx, (icon, title, desc, color) in enumerate(features):
        col = idx % cols
        row = idx // cols
        x = sx + col * (card_w + gap_x)
        y = Inches(1.7) + row * (card_h + gap_y)

        add_shape(slide, x, y, card_w, card_h, BG_CARD, border_color=color, border_width=1)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5),
                     icon, font_size=26)
        add_text_box(slide, x + Inches(0.8), y + Inches(0.2), card_w - Inches(1), Inches(0.35),
                     title, font_size=17, color=color, bold=True)
        add_text_box(slide, x + Inches(0.25), y + Inches(0.75), card_w - Inches(0.5), Inches(1.4),
                     desc, font_size=12, color=LIGHT_GRAY)


def slide_conclusion(prs):
    """Slide 12: Conclusion & Future Work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Top accent
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_PURPLE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "10  Conclusion & Future Work", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3), ACCENT_ROSE)

    # Conclusion card
    add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(3.5), BG_CARD, border_color=ACCENT_BLUE, border_width=1.5)
    add_text_box(slide, Inches(1.1), Inches(1.85), Inches(5), Inches(0.4),
                 "✅ Conclusion", font_size=20, color=ACCENT_BLUE, bold=True)
    add_bullet_list(slide, Inches(1.1), Inches(2.4), Inches(5.3), Inches(2.5), [
        "• Successfully built a hybrid recommendation engine combining three ML approaches",
        "• Collaborative Filtering captures user preference patterns",
        "• SVD extracts meaningful latent features from sparse rating data",
        "• Content-Based Filtering ensures genre relevance",
        "• Full-stack application with polished, responsive UI",
        "• Explainable AI enhances user trust and transparency",
    ], font_size=13, color=LIGHT_GRAY, spacing=Pt(6))

    # Future Work card
    add_shape(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(3.5), BG_CARD, border_color=ACCENT_AMBER, border_width=1.5)
    add_text_box(slide, Inches(7.3), Inches(1.85), Inches(5), Inches(0.4),
                 "🔮 Future Work", font_size=20, color=ACCENT_AMBER, bold=True)
    add_bullet_list(slide, Inches(7.3), Inches(2.4), Inches(5), Inches(2.5), [
        "• Integrate deep learning models (Neural Collaborative Filtering)",
        "• Add user authentication and persistent watch history",
        "• Implement real-time rating feedback loop",
        "• Use TMDB API for richer metadata (cast, directors, trailers)",
        "• Deploy to cloud (AWS / Vercel) for production use",
        "• Add A/B testing framework for algorithm comparison",
    ], font_size=13, color=LIGHT_GRAY, spacing=Pt(6))

    # Thank You section
    add_text_box(slide, Inches(2), Inches(5.8), Inches(9.3), Inches(0.7),
                 "Thank You! 🎬", font_size=38, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(5), Inches(6.5), Inches(3.3), ACCENT_BLUE, Pt(3))
    add_text_box(slide, Inches(2), Inches(6.7), Inches(9.3), Inches(0.5),
                 "Questions & Discussion", font_size=18, color=MUTED_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # Bottom accent
    add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), ACCENT_BLUE)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("🎬 Generating Movie Recommendation System Presentation...")

    slide_title(prs)
    print("  ✓ Slide 1: Title")

    slide_table_of_contents(prs)
    print("  ✓ Slide 2: Table of Contents")

    slide_project_overview(prs)
    print("  ✓ Slide 3: Project Overview")

    slide_architecture(prs)
    print("  ✓ Slide 4: System Architecture")

    slide_algorithms(prs)
    print("  ✓ Slide 5: Recommendation Algorithms")

    slide_tech_stack(prs)
    print("  ✓ Slide 6: Tech Stack")

    slide_frontend_components(prs)
    print("  ✓ Slide 7: Frontend Components")

    slide_api_endpoints(prs)
    print("  ✓ Slide 8: API Endpoints")

    slide_evaluation_metrics(prs)
    print("  ✓ Slide 9: Evaluation Metrics")

    slide_results(prs)
    print("  ✓ Slide 10: Results & Analysis")

    slide_key_features(prs)
    print("  ✓ Slide 11: Key Features")

    slide_conclusion(prs)
    print("  ✓ Slide 12: Conclusion & Future Work")

    output_path = os.path.join("d:/movie-app", OUTPUT_FILE)
    prs.save(output_path)
    print(f"\n🎉 Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
